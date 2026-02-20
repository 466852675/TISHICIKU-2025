import sys
import os

def read_pptx(file_path):
    """使用python-pptx读取PowerPoint文件"""
    if not os.path.exists(file_path):
        print(f"错误: 文件不存在: {file_path}")
        return

    try:
        from pptx import Presentation
    except ImportError:
        print("错误: 未找到python-pptx库，请先安装")
        print("安装命令: pip install python-pptx")
        return

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"错误: 无法打开PowerPoint文件: {file_path}")
        print(f"原因: {e}")
        return

    print(f"=== PowerPoint文档: {file_path} ===")
    print(f"幻灯片数量: {len(prs.slides)}\n")

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*60}")
        print(f"【第 {slide_num} 页】")
        print('='*60)

        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if texts:
            print("\n".join(texts))
        else:
            print("(本页无文本内容)")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python read_pptx.py <文件路径>")
        sys.exit(1)

    read_pptx(sys.argv[1])
